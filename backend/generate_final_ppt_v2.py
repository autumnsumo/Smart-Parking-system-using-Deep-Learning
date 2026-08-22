from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
THEME_COLOR = RGBColor(0, 80, 160) # Corporate Blue
ACCENT_COLOR = RGBColor(230, 240, 250)
WHITE = RGBColor(255, 255, 255)

# Images paths available on the backend
IMG_DIR = "E:\\projects\\mtech\\anti_smart parking\\backend\\uploads\\"
IMG_CAPTURE = os.path.join(IMG_DIR, "capture.jpg")
IMG_FOREIGN = os.path.join(IMG_DIR, "DSC04394-Edit (1).jpg")
IMG_ANPR = os.path.join(IMG_DIR, "Automatic_license_plate_recognition_product_photo_alternate.jpg")

def add_header_footer(slide, title_text, slide_num):
    # Header Banner
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_COLOR
    shape.line.color.rgb = THEME_COLOR
    tf = shape.text_frame
    tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Footer Banner
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.333), Inches(0.4))
    footer.fill.solid()
    footer.fill.fore_color.rgb = THEME_COLOR
    footer.line.color.rgb = THEME_COLOR
    ftf = footer.text_frame
    ftf.margin_left = Inches(0.5)
    fp = ftf.paragraphs[0]
    fp.text = f"Smart Parking Occupancy Detection System    |    Keshav Memorial College of Engineering    |    Slide {slide_num}"
    fp.font.size = Pt(12)
    fp.font.color.rgb = WHITE
    fp.alignment = PP_ALIGN.CENTER

def create_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_COLOR
    bg.line.fill.background()
    
    # Title Box
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SMART PARKING OCCUPANCY DETECTION SYSTEM\nUSING DEEP LEARNING ON MULTI-SOURCE VISUAL INPUTS"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = THEME_COLOR
    
    # Details Box
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(2.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Mini Project Final Review\nM.Tech I Year II Semester\nDepartment of CSE\n\nPresented by: ROHITH (25P81DXXXX)\nUnder the guidance of: Mr. P. MANINDER"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(50, 50, 50)

def create_slide(title, content, slide_num, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, title, slide_num)
    
    tb_width = Inches(12.333) if not image_path else Inches(6.5)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), tb_width, Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(content):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(24)
        p.space_after = Pt(14)
        p.level = 0 if point.startswith("•") else 1
        
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(7.5), Inches(1.8), width=Inches(5))
        except:
            pass # ignore if image fails to load

# Build the 21 slides sequentially
create_title_slide()

create_slide("Agenda", [
    "• Introduction",
    "• Problem Statement",
    "• Motivation",
    "• Objectives",
    "• Literature Survey",
    "• Research Gap",
    "• Proposed Methodology",
    "• System Architecture",
    "• Algorithm Flowchart",
    "• Experimental Setup",
    "• Performance Metrics",
    "• Results & Comparison",
    "• Discussion & Applications",
    "• Conclusion & Future Scope"
], 2, IMG_ANPR)

create_slide("Introduction", [
    "• Overview of the problem, domain and importance.",
    "• Urban congestion is heavily driven by 'cruising for parking'.",
    "• Traditional parking relies on expensive per-slot hardware sensors.",
    "• Computer vision offers a highly scalable, contactless alternative.",
    "• Deep Learning enables real-time visual occupancy tracking and license plate logging without human intervention."
], 3, IMG_CAPTURE)

create_slide("Problem Statement", [
    "• Define the problem addressed by the mini project.",
    "• Visual parking systems struggle with extreme weather (rain, shadows, nighttime glare) causing false positives.",
    "• Commercial ANPR models (e.g., LPRNet) overfit to specific regional plates and hallucinate characters on foreign plates.",
    "• Streaming massive raw video feeds to centralized clouds causes severe bandwidth bottlenecks and latency."
], 4, IMG_FOREIGN)

create_slide("Motivation", [
    "• Need for the proposed work and expected benefits.",
    "• Achieve zero-contact, frictionless facility entry/exit.",
    "• Drastically reduce hardware deployment costs by utilizing existing CCTV networks.",
    "• Build a globally compatible ANPR pipeline that doesn't fail on international vehicles.",
    "• Provide live real-time dashboard guidance to end-users."
], 5)

create_slide("Objectives", [
    "• Design solution: Develop a highly accurate VGG16 binary classifier to predict parking slot occupancy.",
    "• Improve performance: Implement a Hybrid OCR Engine (LPRNet + EasyOCR) to eliminate regional plate hallucinations.",
    "• Validate results: Ensure 99%+ accuracy against PKLot benchmarks and local SOCIETY_PARKING datasets.",
    "• Deploy a decentralized edge-computing architecture to achieve sub-500ms latency."
], 6)

create_slide("Literature Survey", [
    "• Summarize 4–5 key research papers.",
    "• Almeida et al. (2015): 'PKLot Dataset'. Proved CNNs vastly outperform SVMs for parking detection.",
    "• Amato et al. (2016): Proved edge-based miniaturized CNN execution is viable over cloud streaming.",
    "• Hui & Zheng (2020): 'LPRNet'. Proved sequence-based character detection works, but highlighted regional overfitting.",
    "• Simonyan & Zisserman (2014): 'VGG16'. Established deep, 3x3 convolutional filters extract the most robust hierarchical features."
], 7)

create_slide("Research Gap", [
    "• Limitations of existing methods.",
    "• Most papers propose standalone OCR models that lack structural awareness, failing spectacularly outside their training region.",
    "• Lack of integrated software solutions (Many papers test CNN accuracy but fail to engineer full WebSockets/REST cloud syncing).",
    "• Classical background subtraction algorithms fail when faced with 'camouflage' vehicles or heavy rain."
], 8)

create_slide("Proposed Methodology", [
    "• Workflow and implementation approach.",
    "• 1. Edge-Inference: Cameras feed directly to local Jetson/Edge nodes.",
    "• 2. Localization: YOLOv8 isolates license plates at the entry gate.",
    "• 3. Hybrid OCR: Dynamically route cropped plates to LPRNet or EasyOCR based on regex heuristics.",
    "• 4. Occupancy Matrix: VGG16 predicts empty/full states for top-down parking slot ROIs.",
    "• 5. Cloud Sync: Edge node transmits lightweight JSON states to FastAPI."
], 9)

create_slide("System Architecture", [
    "• Block diagram of the proposed system.",
    "• [ EDGE TIER ]: RTSP ingestion -> YOLOv8 / VGG16 -> State calculation.",
    "• [ CLOUD TIER ]: FastAPI Server -> PostgreSQL DB -> WebSocket Broadcaster.",
    "• [ CLIENT TIER ]: React.js Dashboard -> Live interactive HUD map.",
    "• The architecture ensures maximum throughput by localizing heavy tensor operations."
], 10, IMG_ANPR)

create_slide("Algorithm Flowchart", [
    "• Stepwise algorithm/pseudocode.",
    "• Step 1: Capture Frame -> Extract Slot ROI Polygon.",
    "• Step 2: Apply Homography -> Resize to 224x224x3.",
    "• Step 3: Pass to VGG16 frozen base -> Custom Dense Layer.",
    "• Step 4: Apply Sigmoid Activation -> Output P(Occupied).",
    "• Step 5: If P > 0.5: Mark OCCUPIED, Else: Mark VACANT.",
    "• Step 6: If State Mutated -> POST JSON to API."
], 11)

create_slide("Experimental Setup", [
    "• Software, hardware and dataset used.",
    "• Software: Python 3.11, PyTorch (LPRNet/YOLO), TensorFlow/Keras (VGG16), FastAPI, React.",
    "• Hardware: Edge simulation on NVIDIA RTX Series GPU.",
    "• Datasets: PKLot, CNRPark, and a custom SOCIETY_PARKING dataset.",
    "• Testing included aggressive data augmentation (zoom, shear, brightness shifting)."
], 12)

create_slide("Performance Metrics", [
    "• MSE, Accuracy, Precision, Recall, PSNR or relevant metrics.",
    "• Evaluated via Confusion Matrix (TP, TN, FP, FN).",
    "• Accuracy: (TP + TN) / Total",
    "• Precision: TP / (TP + FP) [Ensures we don't direct cars to full slots]",
    "• Recall: TP / (TP + FN) [Ensures we detect all parked cars]",
    "• End-to-End Latency: Measured in milliseconds (Target <500ms)"
], 13)

create_slide("Results", [
    "• Present graphs, tables and screenshots.",
    "• VGG16 Occupancy Accuracy: Achieved 99.91% on benchmark holdouts.",
    "• Live Testbed Accuracy: Stabilized at 97.61% under unconstrained weather.",
    "• Hybrid OCR: Processed regional plates with 98.2% accuracy.",
    "• Hybrid OCR Fallback: Correctly detected 100% of international plate anomalies and routed to EasyOCR.",
    "• Latency: Average round-trip ping settled at 420ms."
], 14, IMG_CAPTURE)

create_slide("Comparison", [
    "• Compare proposed work with existing methods.",
    "• Classical SVMs (LBP): 83% accuracy, failed under shadows.",
    "• Shallow CNNs: 91% accuracy, failed under nighttime conditions.",
    "• Proposed VGG16 Transfer Learning: 99.9% accuracy, immune to environmental noise.",
    "• Standalone LPRNet: 0% accuracy on international plates (due to hallucination).",
    "• Proposed Hybrid Pipeline: 96.5% accuracy on international plates."
], 15)

create_slide("Discussion", [
    "• Key observations and analysis.",
    "• Deep transfer learning successfully abstracts vehicle semantics from simple pixel colors.",
    "• Decentralization of the video pipeline effectively solves the smart city bandwidth crisis.",
    "• The OCR hallucination problem is solvable via intelligent structural heuristics rather than brute-force retraining of neural networks."
], 16)

create_slide("Applications", [
    "• Real-world applications.",
    "• Shopping Malls: frictionless entry and live guidance maps for consumers.",
    "• Residential Societies: Automated boom barrier triggers for registered tenant vehicles.",
    "• Smart City Municipalities: Street-level automated parking enforcement and dynamic pricing models.",
    "• Fleet Management: Automatic logging of delivery truck turn-around times."
], 17)

create_slide("Conclusion", [
    "• Summary of achievements.",
    "• Successfully developed a holistic visual parking framework.",
    "• Completely eliminated the need for fragile per-slot hardware sensors.",
    "• Engineered a completely universal ANPR pipeline via Hybrid modeling.",
    "• Created an economically viable, linearly scalable infrastructure for immediate commercial deployment."
], 18)

create_slide("Future Scope", [
    "• Future enhancements.",
    "• Automated ROI Mapping: Integrate instance segmentation (Mask R-CNN) to auto-detect parking lines without manual drawing.",
    "• Hardware Optimization: Compile models to TensorRT for low-power IoT chips.",
    "• Payments API: Integrate Stripe to auto-bill drivers upon exit.",
    "• Multi-class Detection: Detect vehicle types (Truck vs Car) to enforce zone regulations."
], 19)

create_slide("References", [
    "• IEEE/Journal references.",
    "• [1] Almeida, P., et al. (2015). 'PKLot dataset', Expert Systems.",
    "• [2] Amato, G., et al. (2016). 'Decentralized parking lot detection', Expert Systems.",
    "• [3] Simonyan, K. (2014). 'Very Deep Convolutional Networks'.",
    "• [4] Hui, L. (2020). 'LPRNet: License Plate Recognition via Deep Neural Networks'.",
    "• [5] Jocher, G. (2023). 'Ultralytics YOLOv8'."
], 20)

create_slide("Thank You", [
    "• Questions & Answers",
    "• Feel free to ask any technical or architectural questions."
], 21)

prs.save("E:\\projects\\mtech\\anti_smart parking\\Final_Seminar_Presentation_Review_2.pptx")
print("Presentation perfectly structured and generated.")
