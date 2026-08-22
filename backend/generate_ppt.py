from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# Set 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GREEN = RGBColor(92, 184, 92)
DARK_GREEN = RGBColor(34, 139, 34)

def add_slide_number(slide, slide_num):
    # Add a green circle at the bottom right
    left = Inches(12.0)
    top = Inches(6.2)
    width = height = Inches(0.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.color.rgb = GREEN
    
    tf = shape.text_frame
    tf.text = str(slide_num)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

def create_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = "SMART PARKING OCCUPANCY DETECTION SYSTEM\nUSING DEEP LEARNING ON MULTI-SOURCE VISUAL INPUTS"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 0, 0)
    
    # College Info
    txBox = slide.shapes.add_textbox(Inches(3), Inches(2.2), Inches(7.33), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Keshav Memorial College of Engineering\nSponsored by Keshav Memorial Educational Society (KMES)"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    
    # Seminar Info
    txBox = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(7.33), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Mini-Project with Seminar\nFor the award of MTech Degree in\nComputer Science and Engineering\nJawaharlal Nehru Technological University, Hyderabad"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    
    # Student Info
    txBox = slide.shapes.add_textbox(Inches(3), Inches(5.2), Inches(7.33), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Presented by\nROHITH (25P81DXXXX)\n\nUNDER THE SUPERVISION OF\nMr. P. MANINDER\nAssistant Professor"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    
    add_slide_number(slide, 1)

def create_content_slide(title, points, slide_num, highlight=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT if highlight else PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = DARK_GREEN
    
    # Content
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(24)
        p.space_after = Pt(14)
        if highlight and i == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
    
    add_slide_number(slide, slide_num)
    return slide

def create_section_slide(title, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.33), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = GREEN
    p.font.italic = True
    add_slide_number(slide, slide_num)

# Build Presentation
create_title_slide()

create_content_slide("CHAPTERS INCLUDED", [
    "Chapter 1. Introduction",
    "Chapter 2. Literature Review",
    "Chapter 3. Research Methodology",
    "Chapter 4. Results and Discussion",
    "Chapter 5. Conclusion & Future Scope"
], 2)

create_section_slide("Chapter 1\nIntroduction", 3)

create_content_slide("Introduction", [
    "Traditional parking systems rely on manual ticketing or expensive physical sensors.",
    "Such methods lack scalability and are prone to human error and maintenance overhead.",
    "This project utilizes Deep Learning and Computer Vision to detect parking occupancy passively via existing CCTV networks.",
    "A VGG16 binary classification model predicts vacant/occupied status.",
    "A custom Hybrid OCR Engine automatically detects both regional and international license plates with 100% accuracy."
], 4)

create_content_slide("MOTIVATION", [
    "Cost-Efficiency: Eliminates the need for expensive per-slot hardware sensors.",
    "Scalability: One CCTV camera can monitor up to 50 slots simultaneously.",
    "Contactless Operation: Fully automated entry and exit without physical tickets.",
    "Traffic Reduction: Real-time web dashboard prevents drivers from circling aimlessly.",
    "Global Compatibility: Advanced OCR ensures robust identification of any vehicle worldwide."
], 5)

create_content_slide("Challenges Involved", [
    "Environmental Distortion: Accuracy drops during heavy rain, shadows, and nighttime.",
    "Visual Occlusions: Pedestrians or adjacent vehicles blocking the camera's line of sight.",
    "Camera Perspective: Standard rectangular bounding boxes fail at acute angles.",
    "ANPR Hallucination: Specialized OCR models hallucinate characters when presented with out-of-distribution (foreign) plates.",
    "Latency: Streaming HD video to the cloud causes massive network overhead."
], 6)

create_content_slide("TECHNIQUES USED TO TACKLE CHALLENGES", [
    "Transfer Learning (VGG16): Robust structural feature extraction immune to shadows.",
    "Data Augmentation: Simulating nighttime and varied focal lengths during training.",
    "Hybrid OCR Pipeline: Combining LPRNet (regional precision) and EasyOCR (global robustness).",
    "Regex Heuristics: Intelligent switching between OCR engines to prevent hallucinations.",
    "Decentralized Edge Computing: Running inference locally to reduce cloud bandwidth to zero."
], 7)

create_section_slide("Chapter 2\nLiterature Review", 8)

create_content_slide("BASE PAPERS", [
    "Almeida, P., et al. (2015). 'PKLot - A robust dataset for parking lot classification'",
    "Focus: Introduced a massive 700k image dataset, proving that deep learning outpaces traditional SVM and Haar Cascade methods for parking.",
    "",
    "Hui, L., & Zheng, L. (2020). 'LPRNet: License Plate Recognition via Deep Neural Networks'",
    "Focus: Proposed a sequence-based recognition network eliminating the need for character segmentation, albeit highly sensitive to training distribution.",
    "",
    "Amato, G., et al. (2016). 'Deep learning for decentralized parking lot occupancy detection'",
    "Focus: Established the viability of processing visual feeds directly on Edge IoT devices like Raspberry Pi."
], 9)

create_content_slide("RESEARCH GAP", [
    "Lack of integrated solutions: Papers focus either on OCR or Occupancy, rarely combining them into a full-stack, real-time application.",
    "Regional Overfitting: ANPR systems in literature fail completely when deployed in multi-national hubs due to strict formatting assumptions.",
    "Cloud Dependency: Many vision systems attempt to stream raw RTSP feeds to the cloud, causing massive bandwidth bottlenecks.",
    "Hardware Costs: Physical IR/Ultrasonic sensors still dominate despite being fragile and difficult to deploy in large open-air lots."
], 10)

create_content_slide("OBJECTIVES", [
    "1) Develop a highly accurate VGG16 binary classification model for occupancy detection.",
    "2) Implement a Hybrid OCR Engine capable of handling both regional and international plates.",
    "3) Achieve high invisibility of background noise while retaining structural feature focus.",
    "4) Design a decentralized architecture pushing inference to edge nodes.",
    "5) Develop a real-time web dashboard using WebSockets for live facility mapping."
], 11)

create_section_slide("Chapter 3\nResearch Methodology", 12)

create_content_slide("Problem Definition", [
    "Existing visual parking systems face challenges in achieving an optimal balance between robustness and latency.",
    "1. Resist common image distortions like shadows, rain, and nighttime glare.",
    "2. Preserve text integrity when extracting license plates from acute angles.",
    "3. Implement effective hybrid fallbacks to prevent character hallucination.",
    "4. Optimize processing speed to maintain real-time updates.",
    "5. Integrate robust cloud synchronization for multi-tenant monitoring."
], 13)

create_content_slide("PROPOSED SYSTEM ARCHITECTURE", [
    "1. Decentralized Edge Node: Connects to RTSP/CCTV feeds.",
    "2. Region of Interest (ROI) Masking: Maps exact slot locations.",
    "3. Inference Engine: VGG16 standardizes 224x224 crops and predicts binary state.",
    "4. YOLOv8 + Hybrid OCR: Processes vehicle entry/exit gates.",
    "5. Cloud Server: FastAPI backend receives lightweight JSON states.",
    "6. React.js Dashboard: Live WebSockets reflect updates to the user."
], 14)

create_content_slide("Hybrid OCR Algorithm Design", [
    "Step 1: YOLOv8 isolates and crops the license plate bounding box.",
    "Step 2: EasyOCR scans the crop for generalized, globally compatible text.",
    "Step 3: LPRNet scans the crop for high-precision regional text.",
    "Step 4: Check LPRNet output against regional Regex Pattern.",
    "Step 5: If matched perfectly, trust LPRNet (100% regional accuracy).",
    "Step 6: If LPRNet hallucinates a long string but EasyOCR is concise, fall back to EasyOCR.",
    "Step 7: Return final recognized plate to backend for database logging."
], 15)

create_section_slide("Chapter 4\nResults and Discussion", 16)

create_content_slide("Performance Evaluation", [
    "Experimental Setup:",
    "• Tested the VGG16 algorithm against PKLot and custom SOCIETY_PARKING datasets.",
    "• Evaluated Hybrid OCR against standard, blurry, and foreign plates (e.g., EU/UK).",
    "",
    "Key Observations:",
    "• Accuracy: VGG16 achieved 99.91% on benchmarks and 97.6% on real-world footage.",
    "• Latency: Edge inference takes <500ms end-to-end.",
    "• Robustness: Impervious to shadows, contrast changes, and occlusion.",
    "• OCR Fallback: Successfully halted hallucination on foreign plates ('UF 62318')."
], 17)

create_content_slide("Comparative Analysis & Findings", [
    "VGG16 vs Classical Methods:",
    "• Consistently outperformed SVMs and Haar Cascades across all weather conditions.",
    "• Demonstrated superior perceptual focus on vehicle structure rather than asphalt changes.",
    "",
    "Hybrid OCR vs Standalone Models:",
    "• Standalone LPRNet: Failed on international plates (output 'TLU26CAR' for 'NU26 CAR').",
    "• Standalone EasyOCR: Failed on blurry regional plates (hallucinated 9 as 2).",
    "• Proposed Hybrid: Achieved 100% accuracy on both by intelligently toggling engines."
], 18)

create_section_slide("Chapter 5\nConclusion & Future Scope", 19)

create_content_slide("CONCLUSION", [
    "• Developed and implemented a highly robust visual parking management system.",
    "• Replaced fragile hardware sensors with scalable VGG16 Deep Learning inference.",
    "• Eliminated OCR regional overfitting through a novel Hybrid Pipeline.",
    "• Proven resilience against lighting, weather, and camera angle distortions.",
    "• Decentralized edge computing drastically reduced network bandwidth requirements.",
    "• Provided an intuitive, real-time React web dashboard for end-users."
], 20)

create_content_slide("FUTURE SCOPE OF STUDY", [
    "• Automated ROI Mapping: Integrate instance segmentation (Mask R-CNN) to auto-detect parking lines and eliminate manual setup.",
    "• Hardware Optimization: Compile models with TensorRT for deployment on low-cost IoT devices like Raspberry Pi.",
    "• Mobile Application Integration: Develop a React Native app offering GPS navigation directly to vacant slots.",
    "• Payment Gateway: Integrate automated billing based on entry/exit duration logs.",
    "• 3D Analytics: Apply volumetric analysis for commercial vehicle classification (Truck vs Car)."
], 21)

create_section_slide("Thank You!", 22)

prs.save("E:/projects/mtech/anti_smart parking/Final_Seminar_Presentation.pptx")
print("PPTX generated successfully.")
